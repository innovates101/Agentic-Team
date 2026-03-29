"""
build_presentation.py

Builds a clean, modern PowerPoint presentation on Agentic AI for senior banking executives.
12 slides, 16:9 widescreen, navy/white/gold design, 6 embedded diagrams.
Output: Owner's Inbox/Agentic_AI_Presentation.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Paths ──────────────────────────────────────────────────────────────────────
ROOT = Path(__file__).parent
FIGURES = ROOT / "Owner's Inbox" / "figures"
OUTPUT = ROOT / "Owner's Inbox" / "Agentic_AI_Presentation.pptx"

FIG = {
    1: FIGURES / "fig1_capability_loop.png",
    2: FIGURES / "fig2_genai_vs_agentic.png",
    3: FIGURES / "fig3_autonomy_spectrum.png",
    4: FIGURES / "fig4_tech_stack.png",
    5: FIGURES / "fig5_risk_heatmap.png",
    6: FIGURES / "fig6_strategic_timeline.png",
}

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1B, 0x2A, 0x4A)
NAVY_DARK = RGBColor(0x0F, 0x1A, 0x30)
BLUE      = RGBColor(0x2E, 0x6B, 0xB0)
BLUE_LIGHT= RGBColor(0x5B, 0x9B, 0xD5)
GOLD      = RGBColor(0xC9, 0xA8, 0x4C)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE  = RGBColor(0xF4, 0xF6, 0xFA)
LIGHT_BLU = RGBColor(0xE8, 0xF0, 0xFB)
MID_GRAY  = RGBColor(0x6B, 0x7A, 0x99)
DARK_TEXT = RGBColor(0x1E, 0x1E, 0x2E)
GREEN     = RGBColor(0x1E, 0x7E, 0x5C)
RED       = RGBColor(0xC0, 0x39, 0x2B)

# ── Slide dimensions (16:9) ────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── Helper: create presentation ────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

# ── Helper: blank slide ────────────────────────────────────────────────────────
def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # completely blank

# ── Helper: filled rectangle ──────────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill_color, border=False):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

# ── Helper: add text box ───────────────────────────────────────────────────────
def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb

# ── Helper: add multi-paragraph text box ──────────────────────────────────────
def add_para_text(slide, paragraphs, left, top, width, height,
                  size=16, color=DARK_TEXT, line_spacing=1.2):
    """paragraphs: list of (text, bold, size_override_or_None)"""
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn as _qn
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for (text, bold, size_ov) in paragraphs:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        run = para.add_run()
        run.text = text
        run.font.name = "Calibri"
        run.font.size = _Pt(size_ov if size_ov else size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txb

# ── Helper: standard header bar ───────────────────────────────────────────────
def add_header(slide, title, subtitle=None, bar_color=NAVY):
    # Navy bar across top
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.05), bar_color)
    # Slide title
    add_text(slide, title,
             left=Inches(0.45), top=Inches(0.12), width=Inches(11.5), height=Inches(0.6),
             size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 left=Inches(0.45), top=Inches(0.65), width=Inches(11.5), height=Inches(0.35),
                 size=14, color=BLUE_LIGHT)

# ── Helper: divider slide (section break) ──────────────────────────────────────
def section_divider(prs, number, title, tagline=""):
    slide = blank_slide(prs)
    # Full navy background
    add_rect(slide, 0, 0, W, H, NAVY_DARK)
    # Left accent bar
    add_rect(slide, 0, 0, Inches(0.12), H, GOLD)
    # Section number
    add_text(slide, f"0{number}", left=Inches(0.4), top=Inches(2.2),
             width=Inches(2), height=Inches(1.2), size=72, bold=True, color=GOLD)
    # Title
    add_text(slide, title, left=Inches(0.4), top=Inches(3.3),
             width=Inches(12.5), height=Inches(0.9), size=40, bold=True, color=WHITE)
    if tagline:
        add_text(slide, tagline, left=Inches(0.4), top=Inches(4.2),
                 width=Inches(10), height=Inches(0.5), size=18, color=BLUE_LIGHT, italic=True)
    return slide

# ── Helper: bullet list ───────────────────────────────────────────────────────
def add_bullets(slide, bullets, left, top, width, height,
                size=16, color=DARK_TEXT, marker="▸", marker_color=BLUE, indent=0.3):
    """bullets: list of strings"""
    from pptx.util import Pt as _Pt
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        # Marker run
        r_mark = para.add_run()
        r_mark.text = marker + "  "
        r_mark.font.name = "Calibri"
        r_mark.font.size = _Pt(size)
        r_mark.font.color.rgb = marker_color
        r_mark.font.bold = True
        # Content run
        r_text = para.add_run()
        r_text.text = bullet
        r_text.font.name = "Calibri"
        r_text.font.size = _Pt(size)
        r_text.font.color.rgb = color
        # Paragraph spacing
        pPr = para._p.get_or_add_pPr()
        spcAft = etree.SubElement(pPr, qn("a:spcAft"))
        spcPts = etree.SubElement(spcAft, qn("a:spcPts"))
        spcPts.set("val", "120")  # 12pt space after
    return txb

# ── Helper: stat box ──────────────────────────────────────────────────────────
def add_stat_box(slide, number, label, left, top, width=Inches(2.6), height=Inches(1.5),
                 bg=NAVY, num_color=GOLD, lbl_color=WHITE):
    add_rect(slide, left, top, width, height, bg)
    add_text(slide, number, left+Inches(0.1), top+Inches(0.08), width-Inches(0.2), Inches(0.75),
             size=28, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    add_text(slide, label, left+Inches(0.1), top+Inches(0.75), width-Inches(0.2), Inches(0.65),
             size=11, color=lbl_color, align=PP_ALIGN.CENTER)

# ── Helper: card box ──────────────────────────────────────────────────────────
def add_card(slide, title, body, left, top, width, height,
             title_color=NAVY, body_color=DARK_TEXT, bg=LIGHT_BLU, accent=BLUE):
    add_rect(slide, left, top, width, height, bg)
    add_rect(slide, left, top, Inches(0.06), height, accent)  # left accent
    add_text(slide, title, left+Inches(0.18), top+Inches(0.12), width-Inches(0.3), Inches(0.38),
             size=13, bold=True, color=title_color)
    add_text(slide, body, left+Inches(0.18), top+Inches(0.45), width-Inches(0.3), height-Inches(0.55),
             size=11, color=body_color)


# ══════════════════════════════════════════════════════════════════════════════
#  BUILD SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()

    # ── SLIDE 1: COVER ─────────────────────────────────────────────────────────
    s = blank_slide(prs)
    # Full dark navy background
    add_rect(s, 0, 0, W, H, NAVY_DARK)
    # Gold accent line (bottom third)
    add_rect(s, 0, Inches(5.5), W, Inches(0.06), GOLD)
    # Left accent bar
    add_rect(s, 0, 0, Inches(0.1), H, GOLD)
    # Main title
    add_text(s, "AGENTIC AI",
             Inches(0.5), Inches(1.4), Inches(12.5), Inches(2.0),
             size=80, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # Subtitle
    add_text(s, "The Next Frontier in Banking",
             Inches(0.5), Inches(3.4), Inches(12), Inches(0.8),
             size=32, bold=False, color=GOLD, align=PP_ALIGN.LEFT)
    # Tagline
    add_text(s, "A Strategic Briefing for Senior Banking Executives  |  March 2026",
             Inches(0.5), Inches(4.2), Inches(12), Inches(0.5),
             size=16, color=BLUE_LIGHT, align=PP_ALIGN.LEFT)
    # Bottom note
    add_text(s, "Prepared by the Executive Content & Learning Team",
             Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
             size=12, color=MID_GRAY, align=PP_ALIGN.LEFT)

    # ── SLIDE 2: THE AI EVOLUTION ──────────────────────────────────────────────
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, OFFWHITE)
    add_header(s, "The AI Evolution", "We are at an inflection point")
    add_text(s, "From systems that answer — to systems that act",
             Inches(0.45), Inches(1.2), Inches(12.5), Inches(0.5),
             size=18, italic=True, color=MID_GRAY)

    # Three evolution boxes
    boxes = [
        ("Traditional AI",   "Rule-based automation\nPredefined responses\nNo reasoning",       NAVY,      WHITE, "Pre-2020"),
        ("Generative AI",    "Understands language\nProduces content on demand\nOne prompt → one response", BLUE, WHITE, "2020–2023"),
        ("Agentic AI",       "Receives a goal\nPlans, reasons & acts\nMultiple steps autonomously", GOLD,  NAVY_DARK, "Now →"),
    ]
    box_w = Inches(3.8)
    for i, (title, body, bg, txt, era) in enumerate(boxes):
        lft = Inches(0.45) + i * Inches(4.26)
        add_rect(s, lft, Inches(1.85), box_w, Inches(3.8), bg)
        add_text(s, era, lft + Inches(0.15), Inches(1.95), box_w - Inches(0.3), Inches(0.35),
                 size=11, color=WHITE if bg != GOLD else NAVY, italic=True)
        add_text(s, title, lft + Inches(0.15), Inches(2.3), box_w - Inches(0.3), Inches(0.6),
                 size=20, bold=True, color=txt)
        add_text(s, body, lft + Inches(0.15), Inches(2.95), box_w - Inches(0.3), Inches(2.4),
                 size=15, color=txt)
    # Arrow between boxes (simple text arrows)
    add_text(s, "→", Inches(4.25), Inches(3.5), Inches(0.5), Inches(0.5),
             size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, "→", Inches(8.5), Inches(3.5), Inches(0.5), Inches(0.5),
             size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    # Bottom callout
    add_rect(s, Inches(0.45), Inches(6.1), Inches(12.4), Inches(0.85), NAVY)
    add_text(s, "McKinsey estimates Agentic AI could unlock $200–$340 billion in annual value for global banking  [1]",
             Inches(0.65), Inches(6.2), Inches(12), Inches(0.6),
             size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # ── SLIDE 3: WHAT IS AGENTIC AI? ──────────────────────────────────────────
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, OFFWHITE)
    add_header(s, "What Is Agentic AI?", "An AI that pursues goals — not just answers prompts")

    # Left panel: 4 capabilities
    add_text(s, "Four Core Capabilities",
             Inches(0.45), Inches(1.15), Inches(5.5), Inches(0.45),
             size=14, bold=True, color=NAVY)

    caps = [
        ("🧠  Planning",   "Decomposes a goal into a logical sequence of steps"),
        ("💾  Memory",     "Retains context across steps; accesses institutional knowledge on demand"),
        ("🔧  Tool Use",   "Calls external systems — databases, APIs, calculators, live applications"),
        ("⚡  Action",     "Executes calls, assesses results, and adapts the plan accordingly"),
    ]
    for i, (cap_title, cap_body) in enumerate(caps):
        t = Inches(1.65) + i * Inches(1.28)
        add_rect(s, Inches(0.45), t, Inches(5.5), Inches(1.15), LIGHT_BLU)
        add_rect(s, Inches(0.45), t, Inches(0.07), Inches(1.15), BLUE)
        add_text(s, cap_title, Inches(0.65), t + Inches(0.08), Inches(5.2), Inches(0.38),
                 size=13, bold=True, color=NAVY)
        add_text(s, cap_body, Inches(0.65), t + Inches(0.44), Inches(5.2), Inches(0.62),
                 size=11, color=DARK_TEXT)

    # Right panel: Figure 1
    s.shapes.add_picture(str(FIG[1]), Inches(6.1), Inches(1.1), Inches(6.9), Inches(5.5))

    # Bottom bar
    add_rect(s, Inches(0.45), Inches(6.55), Inches(12.4), Inches(0.7), NAVY)
    add_text(s, "\"Generative AI produces content.  Agentic AI produces results.\"",
             Inches(0.65), Inches(6.62), Inches(12), Inches(0.5),
             size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # ── SLIDE 4: GENAI VS AGENTIC AI ──────────────────────────────────────────
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, OFFWHITE)
    add_header(s, "Generative AI vs Agentic AI", "Same engine — fundamentally different scope")

    # Figure 2 (comparison table) — centred, large
    s.shapes.add_picture(str(FIG[2]), Inches(0.3), Inches(1.1), Inches(9.2), Inches(5.2))

    # Right callout panel
    add_rect(s, Inches(9.65), Inches(1.1), Inches(3.5), Inches(5.2), NAVY)
    add_text(s, "The Critical Distinction",
             Inches(9.82), Inches(1.2), Inches(3.15), Inches(0.45),
             size=13, bold=True, color=GOLD)
    callouts = [
        "Generative AI generates. Agentic AI acts.",
        "Actions carry operational, financial, legal, and regulatory consequences.",
        "A poorly governed agentic system doesn't produce a wrong answer — it executes a wrong transaction.",
        "BCG: The shift to agentic represents the \"single largest near-term productivity unlock\" for banks. [8]",
    ]
    for i, c in enumerate(callouts):
        add_text(s, c, Inches(9.82), Inches(1.75) + i * Inches(1.08),
                 Inches(3.15), Inches(1.0), size=11, color=WHITE)

    # ── SLIDE 5: THREE TYPES ───────────────────────────────────────────────────
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, OFFWHITE)
    add_header(s, "Three Types of Agentic Systems", "Match the type to the risk and regulatory exposure")

    types = [
        (
            "01  Agentic Workflows",
            "Human-in-the-loop at defined decision gates. Process sequence is pre-designed. Predictable and auditable.",
            ["KYC document review", "Loan pre-screening", "Regulatory report generation", "Trade confirmation matching"],
            "Lower risk — the right starting point for most banks in 2026",
            BLUE,
        ),
        (
            "02  Autonomous Agents",
            "Goal-driven. Determines its own plan, selects tools, adapts when results change. Minimal human per step.",
            ["Portfolio monitoring", "Fraud investigation", "Multi-source research synthesis", "Market event response"],
            "High capability ceiling — requires robust guardrails before deployment",
            NAVY,
        ),
        (
            "03  Multi-Agent Systems",
            "Multiple specialist agents collaborating — each handling a step, passing outputs between them at machine speed.",
            ["Deal origination workflows", "End-to-end trade lifecycle", "Complex regulatory reporting", "Cross-function analytics"],
            "The frontier — advanced pilots underway; governance infrastructure still forming [12]",
            GOLD,
        ),
    ]
    card_w = Inches(4.0)
    for i, (title, desc, uses, note, accent_col) in enumerate(types):
        lft = Inches(0.3) + i * Inches(4.35)
        add_rect(s, lft, Inches(1.15), card_w, Inches(5.75), WHITE)
        add_rect(s, lft, Inches(1.15), card_w, Inches(0.08), accent_col)
        add_text(s, title, lft + Inches(0.15), Inches(1.28), card_w - Inches(0.3), Inches(0.55),
                 size=15, bold=True, color=NAVY)
        add_text(s, desc, lft + Inches(0.15), Inches(1.88), card_w - Inches(0.3), Inches(1.2),
                 size=11, color=DARK_TEXT)
        add_text(s, "Banking Applications:",
                 lft + Inches(0.15), Inches(3.12), card_w - Inches(0.3), Inches(0.35),
                 size=11, bold=True, color=accent_col if accent_col != GOLD else NAVY)
        for j, use in enumerate(uses):
            add_text(s, f"• {use}", lft + Inches(0.15), Inches(3.5) + j * Inches(0.35),
                     card_w - Inches(0.3), Inches(0.32), size=11, color=DARK_TEXT)
        # Bottom note
        add_rect(s, lft, Inches(6.4), card_w, Inches(0.5), OFFWHITE)
        add_text(s, note, lft + Inches(0.1), Inches(6.42), card_w - Inches(0.2), Inches(0.45),
                 size=10, italic=True, color=MID_GRAY)

    # ── SLIDE 6: TECHNOLOGY STACK ──────────────────────────────────────────────
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, OFFWHITE)
    add_header(s, "The Technology Stack", "Five layers — understand what you are buying and governing")

    # Figure 4 centred large
    s.shapes.add_picture(str(FIG[4]), Inches(0.25), Inches(1.1), Inches(8.8), Inches(5.7))

    # Right: exec summary of layers
    layers = [
        ("Reasoning Engine", "The LLM — GPT-4o, Claude, Gemini. Cognitive capability of the agent."),
        ("Orchestration",    "Manages planning, execution cycles, and error recovery."),
        ("Tools & APIs",     "Connections to core banking, credit bureaus, market data, compliance systems."),
        ("Memory",           "Short-term (active context) + long-term (institutional knowledge via RAG)."),
        ("Guardrails",       "The delegation of authority matrix for AI. Defines what the agent may do autonomously."),
    ]
    for i, (lyr, desc) in enumerate(layers):
        t = Inches(1.15) + i * Inches(1.1)
        add_rect(s, Inches(9.15), t, Inches(3.85), Inches(1.0), LIGHT_BLU)
        add_rect(s, Inches(9.15), t, Inches(0.07), Inches(1.0), BLUE)
        add_text(s, lyr, Inches(9.33), t + Inches(0.06), Inches(3.5), Inches(0.32),
                 size=12, bold=True, color=NAVY)
        add_text(s, desc, Inches(9.33), t + Inches(0.37), Inches(3.5), Inches(0.56),
                 size=10, color=DARK_TEXT)

    # ── SLIDE 7: BANKING IN ACTION ─────────────────────────────────────────────
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, OFFWHITE)
    add_header(s, "Banking in Action", "Leading institutions are deploying — and measuring results")

    institutions = [
        ("JPMorgan Chase",   "COiN platform processes loan agreements previously requiring 360,000 hrs of annual legal review — now in seconds. LLM Suite deployed to 140,000+ employees. [2]"),
        ("Morgan Stanley",   "AI assistant gives 16,000 financial advisors instant access to 100,000+ proprietary research documents. 'Debrief' auto-drafts meeting follow-ups. [3]"),
        ("Goldman Sachs",    "Agentic coding platform — AI-generated code accounts for a significant and growing share of all new firmwide development. Research agents compress M&A timelines. [4]"),
        ("HSBC",             "AI agents triage AML/fraud alerts, gather evidence from multiple systems, and prioritise cases — materially reducing false-positive rates. [9]"),
        ("DBS Bank",         "800+ AI/ML models in production. Agentic SME lending: loan approvals compressed from days to under one hour. SGD 370M incremental income attributed to AI in 2023. [11]"),
        ("Bank of America",  "Erica virtual assistant: 1.5 billion client interactions. CashPro Assistant for corporate cash management. AI agents automate regulatory reporting workflows. [10]"),
    ]
    cols = 3
    card_w = Inches(4.0)
    card_h = Inches(2.25)
    for i, (name, desc) in enumerate(institutions):
        col = i % cols
        row = i // cols
        lft = Inches(0.35) + col * Inches(4.38)
        t   = Inches(1.15) + row * Inches(2.45)
        add_rect(s, lft, t, card_w, card_h, WHITE)
        add_rect(s, lft, t, card_w, Inches(0.45), NAVY)
        add_text(s, name, lft + Inches(0.12), t + Inches(0.07), card_w - Inches(0.2), Inches(0.32),
                 size=13, bold=True, color=WHITE)
        add_text(s, desc, lft + Inches(0.12), t + Inches(0.52), card_w - Inches(0.2), Inches(1.65),
                 size=10, color=DARK_TEXT)

    # ── SLIDE 8: THE NUMBERS ───────────────────────────────────────────────────
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, NAVY_DARK)
    add_rect(s, 0, 0, Inches(0.1), H, GOLD)
    add_text(s, "The Numbers",
             Inches(0.4), Inches(0.2), Inches(12), Inches(0.7),
             size=30, bold=True, color=WHITE)
    add_text(s, "What the leading institutions have already demonstrated",
             Inches(0.4), Inches(0.82), Inches(12), Inches(0.4),
             size=16, color=BLUE_LIGHT, italic=True)

    # Gold divider
    add_rect(s, Inches(0.4), Inches(1.22), Inches(12.5), Inches(0.04), GOLD)

    stats = [
        ("$200–340B",       "Annual value AI could add to global banking\n(McKinsey, 2023) [1]"),
        ("360,000 hrs",     "Annual legal review compressed to seconds\nJPMorgan COiN [2]"),
        ("SGD 370M",        "Incremental income from AI — one year\nDBS Bank (2023) [11]"),
        ("1.5 Billion",     "Client interactions handled autonomously\nBank of America Erica [10]"),
        ("15%",             "Of daily enterprise decisions made autonomously\nby AI agents by 2028 (Gartner) [7]"),
    ]
    box_w = Inches(2.4)
    box_h = Inches(2.1)
    total_w = len(stats) * box_w + (len(stats) - 1) * Inches(0.15)
    start_l = (W - total_w) / 2
    for i, (num, lbl) in enumerate(stats):
        lft = start_l + i * (box_w + Inches(0.15))
        t   = Inches(1.5)
        add_rect(s, lft, t, box_w, box_h, RGBColor(0x24, 0x3A, 0x60))
        add_rect(s, lft, t, box_w, Inches(0.05), GOLD)
        add_text(s, num, lft + Inches(0.1), t + Inches(0.15), box_w - Inches(0.2), Inches(0.75),
                 size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text(s, lbl, lft + Inches(0.1), t + Inches(0.9), box_w - Inches(0.2), Inches(1.1),
                 size=10, color=WHITE, align=PP_ALIGN.CENTER)

    # Bottom takeaway
    add_rect(s, Inches(0.4), Inches(4.1), Inches(12.5), Inches(0.06), GOLD)
    add_text(s, "30–40% reduction in end-to-end operational costs projected in targeted processes as multi-agent systems scale  (BCG, 2024) [8]",
             Inches(0.4), Inches(4.3), Inches(12.5), Inches(0.55),
             size=14, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

    # Bottom insight
    add_text(s, "Banks deploying agentic systems today are building a compounding advantage that late movers will find difficult to close.",
             Inches(0.4), Inches(5.1), Inches(12.5), Inches(0.6),
             size=14, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

    # ── SLIDE 9: RISKS & GOVERNANCE ───────────────────────────────────────────
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, OFFWHITE)
    add_header(s, "Risks & Governance", "New autonomy demands new control architecture")

    # Left: risk bullets
    risks = [
        ("Cascading Errors",      "Errors propagate and amplify across multi-step processes before detection"),
        ("Accountability Gaps",   "OCC, FCA, and EBA confirm: the institution retains full liability for all AI actions [17,18,19]"),
        ("Prompt Injection",      "Malicious inputs designed to manipulate agent behaviour — a material cybersecurity risk"),
        ("Systemic Correlation",  "Similar agents at multiple banks may create correlated failures in stress scenarios (BIS, FSB) [13,15]"),
        ("Vendor Concentration",  "Reliance on 3–4 LLM providers; unilateral model updates can alter behaviour without notice"),
    ]
    add_text(s, "Five Key Risks",
             Inches(0.45), Inches(1.18), Inches(5.5), Inches(0.4), size=14, bold=True, color=NAVY)
    for i, (risk, desc) in enumerate(risks):
        t = Inches(1.65) + i * Inches(1.0)
        add_rect(s, Inches(0.45), t, Inches(5.5), Inches(0.88), WHITE)
        add_rect(s, Inches(0.45), t, Inches(0.07), Inches(0.88), RED)
        add_text(s, risk, Inches(0.65), t + Inches(0.05), Inches(5.2), Inches(0.3),
                 size=12, bold=True, color=RED)
        add_text(s, desc, Inches(0.65), t + Inches(0.36), Inches(5.2), Inches(0.45),
                 size=10, color=DARK_TEXT)

    # Right: Fig 5 risk heatmap
    s.shapes.add_picture(str(FIG[5]), Inches(6.1), Inches(1.1), Inches(6.9), Inches(5.5))

    # Bottom governance note
    add_rect(s, Inches(0.45), Inches(6.55), Inches(12.4), Inches(0.7), NAVY)
    add_text(s, "Governance imperative: Extend your Model Risk Management framework to cover agentic systems BEFORE deployment — not after.",
             Inches(0.65), Inches(6.62), Inches(12), Inches(0.5),
             size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # ── SLIDE 10: FUTURE OUTLOOK ───────────────────────────────────────────────
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, OFFWHITE)
    add_header(s, "Future Outlook", "Three developments will define the next 24–36 months")

    # Figure 6 (strategic timeline)
    s.shapes.add_picture(str(FIG[6]), Inches(0.25), Inches(1.1), Inches(8.8), Inches(5.7))

    # Right: three horizon bullets
    horizons = [
        ("2026–2027",   "Multi-Agent Systems Scale",
         "Single-agent pilots expand to multi-agent architectures across deal origination, regulatory reporting, and trade lifecycle. BCG projects 30–40% cost reduction in targeted processes. [8]"),
        ("2026–2027",   "Reasoning Models Raise the Bar",
         "Models optimised for careful multi-step reasoning significantly reduce error rates in high-stakes financial decisions — expanding the boundary of what can be deployed autonomously."),
        ("2026–2027",   "Regulatory Clarity Crystallises",
         "EBA guidelines, FCA AI programme, and OCC guidance on agentic model risk expected to mature. Institutions with governance infrastructure already in place convert clarity into competitive advantage. [17,18,19]"),
    ]
    for i, (era, title, body) in enumerate(horizons):
        t = Inches(1.15) + i * Inches(1.85)
        add_rect(s, Inches(9.15), t, Inches(3.95), Inches(1.7), WHITE)
        add_rect(s, Inches(9.15), t, Inches(0.07), Inches(1.7), GOLD)
        add_text(s, era, Inches(9.33), t + Inches(0.06), Inches(3.6), Inches(0.28),
                 size=10, bold=False, color=GOLD, italic=True)
        add_text(s, title, Inches(9.33), t + Inches(0.3), Inches(3.6), Inches(0.38),
                 size=13, bold=True, color=NAVY)
        add_text(s, body, Inches(9.33), t + Inches(0.68), Inches(3.6), Inches(0.9),
                 size=10, color=DARK_TEXT)

    # ── SLIDE 11: RECOMMENDATIONS ─────────────────────────────────────────────
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, OFFWHITE)
    add_header(s, "What Banking Leaders Must Do Now",
               "Act with urgency and governance discipline — the competitive window is open, but not indefinitely")

    columns = [
        (
            "GOVERN",
            NAVY,
            [
                "Extend SR 11-7 model risk management to cover agentic systems — before deployment",
                "Define a Delegation of Authority Matrix for AI agents: what it may do autonomously, what must escalate",
                "Mandate full audit trails for every agent action from day one",
                "Red-team all agents before and after each deployment cycle",
            ],
        ),
        (
            "DEPLOY",
            BLUE,
            [
                "Start with agentic workflows: KYC, document review, regulatory reporting, loan pre-screening",
                "Assign a named AI Agent Owner for every production system — accountability must be human and specific",
                "Reject black-box proposals; insist on explainability and access to full decision logs",
                "Engage regulators proactively — early dialogue positions you as a responsible leader",
            ],
        ),
        (
            "BUILD",
            GREEN,
            [
                "Build agent design and prompt engineering capability internally",
                "Invest in the governance and operating model infrastructure now — before regulatory mandates arrive",
                "Treat Agentic AI as an operating model transformation requiring board-level sponsorship",
                "The institutions combining strategic ambition with governance discipline will set the benchmark",
            ],
        ),
    ]
    col_w = Inches(4.0)
    for i, (title, col, bullets) in enumerate(columns):
        lft = Inches(0.35) + i * Inches(4.38)
        # Header
        add_rect(s, lft, Inches(1.15), col_w, Inches(0.6), col)
        add_text(s, title, lft + Inches(0.12), Inches(1.22), col_w - Inches(0.2), Inches(0.45),
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Body
        add_rect(s, lft, Inches(1.75), col_w, Inches(5.15), WHITE)
        for j, b in enumerate(bullets):
            t = Inches(1.85) + j * Inches(1.22)
            add_rect(s, lft + Inches(0.12), t, Inches(0.06), Inches(0.95), col)
            add_text(s, b, lft + Inches(0.3), t + Inches(0.05), col_w - Inches(0.45), Inches(0.9),
                     size=11, color=DARK_TEXT)

    # ── SLIDE 12: CLOSING ──────────────────────────────────────────────────────
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, NAVY_DARK)
    add_rect(s, 0, 0, Inches(0.1), H, GOLD)
    add_rect(s, 0, Inches(4.9), W, Inches(0.06), GOLD)

    # Quote
    add_text(s,
             "\"The institutions that combine strategic ambition\n"
             "with governance discipline will set the competitive benchmark.\n"
             "Those waiting for perfect certainty will find themselves\n"
             "playing catch-up in a game where the leaders are\n"
             "already several moves ahead.\"",
             Inches(0.6), Inches(0.9), Inches(12.2), Inches(3.7),
             size=28, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

    add_text(s, "Act now. The window is open.",
             Inches(0.6), Inches(4.3), Inches(12), Inches(0.7),
             size=34, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    add_text(s, "Questions?  |  Further reading and references in the accompanying deep-dive report.",
             Inches(0.6), Inches(5.2), Inches(12), Inches(0.5),
             size=15, color=BLUE_LIGHT, align=PP_ALIGN.LEFT)

    add_text(s, "Prepared by the Executive Content & Learning Team  |  March 2026",
             Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
             size=12, color=MID_GRAY, align=PP_ALIGN.LEFT)

    # ── SAVE ───────────────────────────────────────────────────────────────────
    prs.save(str(OUTPUT))
    size_kb = OUTPUT.stat().st_size // 1024
    print(f"Saved:  {OUTPUT}")
    print(f"Size:   {size_kb:,} KB")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
